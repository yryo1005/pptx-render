"""回帰テスト用のPPTXフィクスチャ（tests/fixtures/*.pptx）を生成するスクリプト．

`python-pptx`を用いて，以下の10種類の最小構成PPTXを生成する．

1. empty_slide  : 空スライド
2. text         : テキストボックス（書式混在）
3. japanese     : 日本語テキスト（折り返し確認用）
4. english      : 英語テキスト（折り返し確認用）
5. image        : 画像（PNG）
6. shapes       : 図形（四角形・円・線・矢印・多角形）
7. multi        : 複数要素（テキスト＋画像＋図形）
8. table        : 表（セル結合を含む）
9. equation     : 数式（OMML）
10. complex     : 複合スライド（上記を組み合わせたもの）

生成されたPPTXは，`tests/test_render.py`から回帰テストの入力として使用する．
"""

from __future__ import annotations

import io
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

FIXTURES_DIR = Path(__file__).parent / "fixtures"

_M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

_OMATH_FRACTION_XML = f"""
<a14:m xmlns:a14="{_A14_NS}">
  <m:oMath xmlns:m="{_M_NS}">
    <m:f>
      <m:num><m:r><m:t>1</m:t></m:r></m:num>
      <m:den><m:r><m:t>N</m:t></m:r></m:den>
    </m:f>
    <m:r><m:t>&#8721;</m:t></m:r>
    <m:sSub>
      <m:e><m:r><m:t>x</m:t></m:r></m:e>
      <m:sub><m:r><m:t>i</m:t></m:r></m:sub>
    </m:sSub>
  </m:oMath>
</a14:m>
"""


def _make_png_bytes(color: tuple[int, int, int], size=(200, 120)) -> bytes:
    img = Image.new("RGB", size, color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def gen_empty_slide() -> None:
    prs = Presentation()
    _add_blank_slide(prs)
    prs.save(FIXTURES_DIR / "01_empty_slide.pptx")


def gen_text() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = "Bold and "
    r0.font.bold = True
    r0.font.size = Pt(24)
    r1 = p0.add_run()
    r1.text = "Italic and "
    r1.font.italic = True
    r1.font.size = Pt(24)
    r2 = p0.add_run()
    r2.text = "Underline"
    r2.font.underline = True
    r2.font.size = Pt(24)

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r3 = p1.add_run()
    r3.text = "Centered colored text"
    r3.font.size = Pt(20)
    r3.font.color.rgb = RGBColor(0xCC, 0x22, 0x22)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r4 = p2.add_run()
    r4.text = "Right aligned"
    r4.font.size = Pt(20)

    prs.save(FIXTURES_DIR / "02_text.pptx")


def gen_japanese() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = (
        "日本語のテキストは，行の折り返し位置がPowerPointと大きく異ならないことが"
        "重要である。句読点や括弧などの禁則処理も確認する。（括弧の例）「かぎ括弧の例」"
    )
    r0.font.size = Pt(18)
    prs.save(FIXTURES_DIR / "03_japanese.pptx")


def gen_english() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = (
        "This is a long English paragraph used to verify that word wrapping "
        "behaves correctly when the text box width is limited."
    )
    r0.font.size = Pt(18)
    prs.save(FIXTURES_DIR / "04_english.pptx")


def gen_image() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)
    png_bytes = _make_png_bytes((66, 133, 244))
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(1), Inches(1), width=Inches(4), height=Inches(2.4))
    prs.save(FIXTURES_DIR / "05_image.pptx")


def gen_shapes() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    rect.line.color.rgb = RGBColor(0x00, 0x00, 0x00)

    rr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(0.5), Inches(2), Inches(1))
    rr.fill.solid()
    rr.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)

    ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(0.5), Inches(2), Inches(1))
    ellipse.fill.solid()
    ellipse.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)
    ellipse.fill.transparency = 0.0

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5), Inches(2), Inches(2), Inches(1))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)

    triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3), Inches(2), Inches(1.5), Inches(1.2))
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(0x5B, 0x9B, 0xD5)

    star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(5.5), Inches(2), Inches(1.5), Inches(1.5))
    star.fill.solid()
    star.fill.fore_color.rgb = RGBColor(0xA5, 0x2A, 0x2A)

    line = slide.shapes.add_connector(1, Inches(0.5), Inches(4), Inches(3), Inches(5))
    line.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    line.line.width = Pt(2)

    prs.save(FIXTURES_DIR / "06_shapes.pptx")


def gen_multi() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "複数要素を含むスライド"
    r.font.size = Pt(28)
    r.font.bold = True

    png_bytes = _make_png_bytes((100, 200, 100))
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.5), Inches(1.5), width=Inches(3), height=Inches(1.8))

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.5), Inches(2), Inches(1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

    prs.save(FIXTURES_DIR / "07_multi.pptx")


def gen_table() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)
    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(8), Inches(2))
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = f"R{r}C{c}"
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

    table.cell(1, 0).merge(table.cell(1, 1))
    prs.save(FIXTURES_DIR / "08_table.pptx")


def gen_equation() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1.5))
    tf = box.text_frame
    p_el = tf.paragraphs[0]._p
    m_el = etree.fromstring(_OMATH_FRACTION_XML.encode("utf-8"))
    p_el.append(m_el)
    prs.save(FIXTURES_DIR / "09_equation.pptx")


def gen_complex() -> None:
    prs = Presentation()
    slide = _add_blank_slide(prs)

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.8))
    r = title.text_frame.paragraphs[0].add_run()
    r.text = "複合スライド（Complex Slide）"
    r.font.size = Pt(28)
    r.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(4.5), Inches(1.5))
    r2 = body.text_frame.paragraphs[0].add_run()
    r2.text = "日本語とEnglishが混在する本文．数式・図・表も含む．"
    r2.font.size = Pt(16)

    box = slide.shapes.add_textbox(Inches(0.4), Inches(2.6), Inches(4.5), Inches(1.0))
    p_el = box.text_frame.paragraphs[0]._p
    m_el = etree.fromstring(_OMATH_FRACTION_XML.encode("utf-8"))
    p_el.append(m_el)

    png_bytes = _make_png_bytes((200, 150, 50))
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(5.2), Inches(1.0), width=Inches(3.5), height=Inches(2.1))

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.3), Inches(3.5), Inches(0.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)
    rt = rect.text_frame.paragraphs[0].add_run()
    rt.text = "Rounded rect with text"
    rt.font.size = Pt(14)

    table_shape = slide.shapes.add_table(2, 2, Inches(0.4), Inches(3.8), Inches(4.5), Inches(1.2))
    table = table_shape.table
    for r_idx in range(2):
        for c_idx in range(2):
            table.cell(r_idx, c_idx).text = f"cell {r_idx}-{c_idx}"

    group_shapes = slide.shapes.add_group_shape(
        [
            slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.2), Inches(1), Inches(1)),
        ]
    )

    prs.save(FIXTURES_DIR / "10_complex.pptx")


def main() -> None:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    gen_empty_slide()
    gen_text()
    gen_japanese()
    gen_english()
    gen_image()
    gen_shapes()
    gen_multi()
    gen_table()
    gen_equation()
    gen_complex()
    print(f"生成完了: {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
